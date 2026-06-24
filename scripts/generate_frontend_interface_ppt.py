from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "artifacts" / "outputs"
OUTPUT_FILE = OUTPUT_DIR / "bosch-frontend-interface-jd-review.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

IMG = ROOT / "artifacts" / "outputs" / "ui-captures"
LOGIN = IMG / "login-current.png"
ADMIN_HOME = IMG / "admin-home-current.png"
ANALYTICS = IMG / "analytics-current.png"
OPERATIONS = IMG / "operations-current.png"
PREDICTIVE = IMG / "predictive-current.png"
ACCOUNTS = IMG / "accounts-current.png"
AUDIT = IMG / "audit-current.png"
VISUAL = ROOT / "assets" / "auth-visual-bosch-family.jpg"

BOSCH_RED = RGBColor(0xE2, 0x00, 0x15)
BOSCH_RED_DARK = RGBColor(0xA9, 0x00, 0x11)
BOSCH_BLUE = RGBColor(0x00, 0x56, 0x91)
CHARCOAL = RGBColor(0x1A, 0x1C, 0x1C)
MUTED = RGBColor(0x5E, 0x66, 0x73)
LINE = RGBColor(0xD6, 0xD9, 0xDE)
BG = RGBColor(0xF3, 0xF3, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_RED = RGBColor(0xFF, 0xEB, 0xE9)
SOFT_BLUE = RGBColor(0xE8, 0xF2, 0xFB)
SOFT_GRAY = RGBColor(0xF8, 0xF8, 0xF8)
SUCCESS = RGBColor(0x15, 0x80, 0x3D)
WARNING = RGBColor(0xD9, 0x77, 0x06)

FONT_HEAD = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"


def hex_color(color):
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_lines(slide, lines, left, top, width, height, size=14, color=CHARCOAL):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = item
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_top(slide, label="Bosch Power Tools · Frontend Prototype", page=None, dark=False):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BOSCH_RED
    bar.line.fill.background()

    color = WHITE if dark else BOSCH_RED
    add_text(slide, label, 0.48, 0.28, 4.8, 0.28, size=9.5, color=color, bold=True)
    if page:
        add_text(slide, f"{page:02d}", 12.33, 7.18, 0.45, 0.18, size=9.5, color=WHITE if dark else MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, dark=False):
    add_text(
        slide,
        "Bosch Power Tools | YuYu | 前端界面方向评审",
        0.48,
        7.18,
        7.0,
        0.18,
        size=8.5,
        color=WHITE if dark else MUTED,
    )


def add_title(slide, kicker, title, subtitle="", page=None):
    add_top(slide, page=page)
    add_text(slide, kicker.upper(), 0.72, 0.72, 4.8, 0.25, size=9.5, color=BOSCH_RED, bold=True)
    add_text(slide, title, 0.72, 1.03, 6.4, 0.9, size=29, color=CHARCOAL, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.9, 7.2, 0.38, size=12.5, color=MUTED)
    add_footer(slide)


def add_pill(slide, text, left, top, width, fill, color=CHARCOAL):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = color
    return shape


def add_card(slide, title, body, left, top, width, height, accent=BOSCH_RED, fill=WHITE, tag=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.shadow.inherit = False

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    if tag:
        add_pill(slide, tag, left + 0.22, top + 0.18, 1.05, SOFT_BLUE if accent == BOSCH_BLUE else SOFT_RED, accent)
        title_y = top + 0.64
    else:
        title_y = top + 0.22
    add_text(slide, title, left + 0.22, title_y, width - 0.44, 0.34, size=15.5, color=CHARCOAL, bold=True)
    add_rich_lines(slide, body.split("\n"), left + 0.22, title_y + 0.48, width - 0.44, height - 0.78, size=11.5, color=MUTED)
    return shape


def add_number_card(slide, number, title, body, left, top, width, height, accent=BOSCH_RED):
    add_card(slide, title, body, left, top, width, height, accent=accent, fill=WHITE)
    add_text(slide, number, left + width - 0.72, top + 0.14, 0.45, 0.32, size=17, color=accent, bold=True, align=PP_ALIGN.RIGHT)


def image_size(path):
    with Image.open(path) as img:
        return img.size


def add_image(slide, path, left, top, width, height, caption=None, border=True):
    if border:
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = LINE
    pad = 0.1 if border else 0
    cap_h = 0.32 if caption else 0
    box_left, box_top = left + pad, top + pad
    box_w, box_h = width - pad * 2, height - pad * 2 - cap_h
    img_w, img_h = image_size(path)
    ratio = min(box_w / img_w, box_h / img_h)
    w, h = img_w * ratio, img_h * ratio
    x = box_left + (box_w - w) / 2
    y = box_top + (box_h - h) / 2
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if caption:
        add_text(slide, caption, left + 0.15, top + height - 0.3, width - 0.3, 0.18, size=8.5, color=MUTED)


def add_flow(slide, items, left, top, width, accent=BOSCH_RED):
    gap = 0.18
    item_w = (width - gap * (len(items) - 1)) / len(items)
    for idx, item in enumerate(items):
        x = left + idx * (item_w + gap)
        fill = SOFT_RED if idx == 0 else SOFT_BLUE if idx == 1 else SOFT_GRAY
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(item_w), Inches(0.78))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = LINE
        add_text(slide, item, x + 0.1, top + 0.19, item_w - 0.2, 0.28, size=11.5, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(items) - 1:
            add_text(slide, ">", x + item_w + 0.04, top + 0.22, 0.1, 0.25, size=14, color=accent, bold=True, align=PP_ALIGN.CENTER)


def add_mini_nav(slide, items, left, top, width, height):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    add_text(slide, "页面模块", left + 0.2, top + 0.18, width - 0.4, 0.25, size=12.5, color=BOSCH_RED, bold=True)
    for idx, item in enumerate(items):
        y = top + 0.62 + idx * 0.38
        fill = SOFT_RED if idx == 0 else SOFT_GRAY
        dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.18), Inches(y), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BOSCH_RED if idx == 0 else LINE
        dot.line.fill.background()
        add_text(slide, item, left + 0.4, y - 0.04, width - 0.55, 0.2, size=10.3, color=CHARCOAL if idx < 6 else MUTED)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.author = "YuYu"
    prs.core_properties.title = "设备管理前端界面原型阶段汇报"
    prs.core_properties.subject = "基于实习岗位 JD 的前端页面方向评审"

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CHARCOAL)
    add_top(slide, dark=True, page=1)
    cover_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.14), Inches(4.65), Inches(7.36))
    cover_panel.fill.solid()
    cover_panel.fill.fore_color.rgb = CHARCOAL
    cover_panel.line.fill.background()
    add_text(slide, "FRONTEND PROTOTYPE", 0.72, 0.9, 3.2, 0.25, size=9.5, color=WHITE, bold=True)
    add_text(slide, "设备管理\n前端界面原型\n阶段汇报", 0.72, 1.36, 4.5, 1.85, size=32, color=WHITE, bold=True)
    add_text(slide, "基于实习岗位 JD 先设计的页面方向假设", 0.74, 3.45, 3.7, 0.38, size=13.5, color=RGBColor(0xEA, 0xEA, 0xEA))
    add_pill(slide, "本次只展示前端界面", 0.74, 4.18, 2.0, SOFT_RED, BOSCH_RED_DARK)
    add_rich_lines(slide, ["场景：设备管理数字化探索", "汇报人：YuYu", "日期：2026.04"], 0.76, 4.75, 3.3, 0.95, size=12, color=RGBColor(0xEA, 0xEA, 0xEA))
    if LOGIN.exists():
        add_image(slide, LOGIN, 4.88, 0.68, 7.95, 5.93, caption="当前登录页前端界面", border=True)
    add_footer(slide, dark=True)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "Stage 1 / Background", "为什么我先做这个前端界面", "先根据实习岗位 JD 把自己的理解可视化，便于和 Boss 快速对齐。", page=2)
    add_number_card(slide, "01", "输入来自 JD", "我先根据岗位描述里对设备管理、数据分析和数字化能力的要求，推测这个岗位可能需要支持的页面方向。", 0.72, 2.45, 3.75, 1.35, accent=BOSCH_RED)
    add_number_card(slide, "02", "真实业务还没完整了解", "目前还没有完整掌握产线、设备类型、设备清单、字段口径和现场数据。", 4.82, 2.45, 3.55, 1.35, accent=BOSCH_BLUE)
    add_number_card(slide, "03", "先做前端假设", "所以我先做一版前端界面，把我理解中的设备管理页面结构展示出来。", 8.72, 2.45, 3.85, 1.35, accent=WARNING)
    add_flow(slide, ["实习岗位 JD", "我的理解", "前端界面", "Boss 反馈", "页面调整"], 0.84, 4.45, 11.6, accent=BOSCH_RED)
    add_card(slide, "这一页要讲清楚的核心", "这不是最终方案，也不是完整系统；它是一个页面方向假设，用来确认我的理解是否和 Boss 的期待一致。", 0.84, 5.7, 11.6, 0.86, accent=BOSCH_RED, fill=WHITE)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "Stage 1 / Boundary", "这次只展示前端页面，不讲后端和真实交互", "避免把演示重点拉偏：这次先看页面方向、信息层级和视觉结构。", page=3)
    add_card(slide, "这次展示什么", "登录页视觉\n顶部全局区域\n左侧页面导航\n管理员页面框架\n普通账号页面想法", 0.72, 2.35, 3.85, 3.4, accent=BOSCH_RED, fill=SOFT_RED)
    add_card(slide, "这次不展示什么", "后端接口\n数据库设计\n真实登录\n真实提交审核\n正式业务结论", 4.86, 2.35, 3.85, 3.4, accent=BOSCH_BLUE, fill=SOFT_BLUE)
    add_card(slide, "希望 Boss 看什么", "页面方向是否符合期待\n模块命名是否合理\n首页信息优先级是否正确\n普通账号页面是否有必要\n下一版先改哪里", 9.0, 2.35, 3.35, 3.4, accent=SUCCESS, fill=SOFT_GRAY)
    add_pill(slide, "关键词：前端界面方向评审", 4.55, 6.2, 3.9, SOFT_RED, BOSCH_RED_DARK)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "Stage 2 / Framework", "当前界面先做出了统一入口和页面框架", "登录页、顶部区域、左侧导航和内容工作区先搭出来，用来承载后续页面展示。", page=4)
    if LOGIN.exists():
        add_image(slide, LOGIN, 0.72, 2.25, 5.25, 3.95, caption="登录页：品牌入口、账号区域、基础视觉风格", border=True)
    if ADMIN_HOME.exists():
        add_image(slide, ADMIN_HOME, 6.2, 2.25, 6.25, 3.95, caption="管理员首页：顶部全局区、左侧导航、内容工作区", border=True)
    add_flow(slide, ["登录页", "顶部全局区", "左侧导航", "内容展示区"], 1.1, 6.08, 10.75, accent=BOSCH_RED)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "Stage 2 / Admin Pages", "管理员端先展示总览、分析、设备和规则页面", "这里重点讲页面为什么这样排，不讲功能是否已经真正连通。", page=5)
    add_mini_nav(slide, ["数据提报页面", "生命周期", "趋势分析", "设备看板", "布局示意", "风险规则", "预测性维护", "账户权限页面", "审计日志页面"], 0.72, 2.15, 2.55, 4.55)
    if ADMIN_HOME.exists():
        add_image(slide, ADMIN_HOME, 3.55, 2.15, 3.05, 2.16, caption="首页：KPI 与整体排版", border=True)
    if ANALYTICS.exists():
        add_image(slide, ANALYTICS, 6.82, 2.15, 2.9, 2.16, caption="趋势分析：图表和指标", border=True)
    if OPERATIONS.exists():
        add_image(slide, OPERATIONS, 9.94, 2.15, 2.75, 2.16, caption="设备看板：设备卡片", border=True)
    add_card(slide, "我对管理员页面的假设", "Boss 可能希望先看到一个管理视角页面框架：先看总览，再看趋势和设备，再看风险规则或扩展页面。", 3.55, 4.72, 4.25, 1.35, accent=BOSCH_RED, fill=SOFT_RED)
    add_card(slide, "本页不强调", "真实登录、真实审核、数据落库、权限流转等功能；这里只看前端页面结构是否合理。", 8.08, 4.72, 4.61, 1.35, accent=BOSCH_BLUE, fill=SOFT_BLUE)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "Stage 2 / Device Story", "围绕 OEE 和风险，先设计一条页面展示主线", "我的假设：管理层先看整体状态，再快速看到重点设备。", page=6)
    if ANALYTICS.exists():
        add_image(slide, ANALYTICS, 0.72, 2.15, 4.05, 2.85, caption="趋势分析页：OEE、T-loss、MTBF、MTTR 的呈现方式", border=True)
    if OPERATIONS.exists():
        add_image(slide, OPERATIONS, 4.98, 2.15, 4.05, 2.85, caption="设备看板页：设备卡片、高风险项、关键备件区域", border=True)
    if PREDICTIVE.exists():
        add_image(slide, PREDICTIVE, 9.24, 2.15, 3.28, 2.85, caption="预测性维护页：未来扩展页面的视觉方向", border=True)
    add_flow(slide, ["看全局指标", "识别趋势变化", "定位高风险设备", "查看规则展示"], 0.92, 5.55, 11.1, accent=BOSCH_RED)
    add_text(slide, "请 Boss 判断：这种从首页到趋势、设备、风险的页面主线，是否符合管理层实际想看的顺序。", 1.2, 6.55, 10.8, 0.35, size=13.5, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "Stage 2 / Ordinary Account", "普通账号侧只是先展示一个页面假设", "如果未来需要现场人员查看或填写信息，界面可以这样组织。", page=7)
    if ACCOUNTS.exists():
        add_image(slide, ACCOUNTS, 0.72, 2.15, 5.2, 3.45, caption="普通账号相关页面样式参考：表单、权限、列表区域", border=True)
    if AUDIT.exists():
        add_image(slide, AUDIT, 6.15, 2.15, 5.95, 3.45, caption="记录列表样式参考：时间、事件、对象的展示方式", border=True)
    add_flow(slide, ["查看设备状态", "填写变化信息", "查看记录列表"], 1.25, 5.82, 10.65, accent=BOSCH_BLUE)
    callout = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.28), Inches(6.65), Inches(10.6), Inches(0.5))
    callout.fill.solid()
    callout.fill.fore_color.rgb = SOFT_RED
    callout.line.color.rgb = LINE
    add_text(
        slide,
        "本页需要 Boss 判断：普通账号侧是否需要单独设计页面，还是只展示管理员视角就够了。",
        1.48,
        6.78,
        10.2,
        0.22,
        size=12,
        color=CHARCOAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CHARCOAL)
    add_top(slide, dark=True, page=8)
    add_text(slide, "STAGE 3 / BOSS FEEDBACK", 0.72, 0.86, 4.5, 0.25, size=9.5, color=WHITE, bold=True)
    add_text(slide, "我希望这次汇报后\n先确认前端方向", 0.72, 1.25, 5.1, 1.25, size=28, color=WHITE, bold=True)
    add_text(slide, "如果方向一致，下一步先调页面结构和展示重点；如果方向有偏差，现在改页面成本最低。", 0.74, 2.72, 5.2, 0.55, size=13, color=RGBColor(0xEA, 0xEA, 0xEA))
    add_card(slide, "请 Boss 重点判断", "1. 页面方向是否符合岗位期待\n2. 首页指标优先级是否正确\n3. 9 个页面入口是否太多或太少\n4. 普通账号页面是否有必要\n5. 下一版先改结构还是补内容", 6.2, 0.92, 5.9, 2.6, accent=BOSCH_RED, fill=WHITE)
    add_card(slide, "下一步建议", "1. 先根据反馈调整首页和导航\n2. 合并或删减不必要页面\n3. 强化设备看板和趋势分析\n4. 再考虑补真实设备截图或样例数据", 6.2, 3.82, 5.9, 2.05, accent=BOSCH_BLUE, fill=WHITE)
    add_flow(slide, ["方向确认", "页面调整", "模块取舍", "第二版前端界面"], 0.74, 4.3, 5.2, accent=BOSCH_RED)
    add_footer(slide, dark=True)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    return OUTPUT_FILE


if __name__ == "__main__":
    print(build())
