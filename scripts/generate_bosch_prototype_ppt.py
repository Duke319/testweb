from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "artifacts" / "outputs"
OUTPUT_FILE = OUTPUT_DIR / "bosch-prototype-progress-yuyu-skills.pptx"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

AUTH_SCREEN = ROOT / "artifacts" / "outputs" / "ui-captures" / "login-current.png"
ADMIN_SCREEN = ROOT / "artifacts" / "outputs" / "ui-captures" / "admin-home-current.png"
ANALYTICS_SCREEN = ROOT / "artifacts" / "outputs" / "ui-captures" / "analytics-current.png"
EQUIPMENT_SCREEN = ROOT / "artifacts" / "outputs" / "ui-captures" / "operations-current.png"
PREDICTIVE_SCREEN = ROOT / "artifacts" / "outputs" / "ui-captures" / "predictive-current.png"
ACCOUNTS_SCREEN = ROOT / "artifacts" / "outputs" / "ui-captures" / "accounts-current.png"
AUDIT_SCREEN = ROOT / "artifacts" / "outputs" / "ui-captures" / "audit-current.png"
AUTH_VISUAL = ROOT / "assets" / "auth-visual-bosch-family.jpg"


BOSCH_RED = RGBColor(0xE2, 0x00, 0x15)
BOSCH_RED_DARK = RGBColor(0xB4, 0x00, 0x0E)
BOSCH_BLUE = RGBColor(0x00, 0x56, 0x91)
TEXT = RGBColor(0x1A, 0x1C, 0x1C)
MUTED = RGBColor(0x5B, 0x63, 0x70)
LIGHT_BG = RGBColor(0xF3, 0xF3, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_RED = RGBColor(0xFF, 0xEB, 0xE9)
SOFT_BLUE = RGBColor(0xE8, 0xF2, 0xFB)
SOFT_GRAY = RGBColor(0xF8, 0xF8, 0xF8)
LINE = RGBColor(0xD7, 0xD9, 0xDE)
SUCCESS = RGBColor(0x15, 0x80, 0x3D)
WARNING = RGBColor(0xD9, 0x77, 0x06)

FONT_HEAD = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_topbar(slide, title="Bosch Power Tools", slide_num=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.17))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BOSCH_RED
    bar.line.fill.background()

    brand = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(2.8), Inches(0.3))
    tf = brand.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_BODY
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = BOSCH_RED

    if slide_num is not None:
        num_box = slide.shapes.add_textbox(SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.38), Inches(0.4), Inches(0.2))
        tf = num_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{slide_num}"
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.color.rgb = MUTED


def add_footer(slide, text="Bosch Power Tools | YuYu | Prototype Progress Review"):
    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.0), SLIDE_W - Inches(0.9), Inches(0.22))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_title_block(slide, kicker, title, subtitle="", left=0.6, top=0.6, width=5.4):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = kicker.upper()
    run.font.name = FONT_BODY
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = BOSCH_RED

    p = tf.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEAD
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TEXT

    if subtitle:
        p = tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = subtitle
        run.font.name = FONT_BODY
        run.font.size = Pt(13)
        run.font.color.rgb = MUTED


def add_chip(slide, text, left, top, width, fill_rgb, text_rgb=TEXT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = text_rgb


def add_bullet_box(slide, heading, bullets, left, top, width, height, fill=WHITE, accent=BOSCH_RED):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = LINE

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    title = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.14), Inches(width - 0.28), Inches(0.32))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = heading
    run.font.name = FONT_HEAD
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT

    body = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.58), Inches(width - 0.32), Inches(height - 0.7))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(5)
        p.bullet = True
        run = p.add_run()
        run.text = bullet
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.color.rgb = TEXT


def add_card(slide, title, body, left, top, width, height, tag=None, accent=BOSCH_RED, fill=WHITE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    if tag:
        add_chip(slide, tag, left + 0.2, top + 0.18, 1.1, SOFT_BLUE if accent == BOSCH_BLUE else SOFT_RED, accent)

    title_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.55), Inches(width - 0.4), Inches(0.5))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEAD
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT

    body_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.0), Inches(width - 0.4), Inches(height - 1.15))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.name = FONT_BODY
    run.font.size = Pt(13)
    run.font.color.rgb = MUTED


def add_text_list(slide, title, items, left, top, width, height):
    title_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.32))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEAD
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT

    body = slide.shapes.add_textbox(Inches(left), Inches(top + 0.38), Inches(width), Inches(height - 0.38))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(5)
        p.bullet = True
        run = p.add_run()
        run.text = item
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.color.rgb = TEXT


def fit_image_in_box(path, box_w, box_h):
    with Image.open(path) as img:
        img_w, img_h = img.size
    ratio = min(box_w / img_w, box_h / img_h)
    return img_w * ratio, img_h * ratio


def add_image_panel(slide, image_path, caption, left, top, width, height):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = LINE

    img_left = Inches(left + 0.12)
    img_top = Inches(top + 0.12)
    box_w = Inches(width - 0.24)
    box_h = Inches(height - 0.52)

    fit_w, fit_h = fit_image_in_box(image_path, box_w, box_h)
    pic_left = img_left + (box_w - fit_w) / 2
    pic_top = img_top + (box_h - fit_h) / 2
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=fit_w, height=fit_h)

    cap = slide.shapes.add_textbox(Inches(left + 0.14), Inches(top + height - 0.32), Inches(width - 0.28), Inches(0.18))
    tf = cap.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = caption
    run.font.name = FONT_BODY
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def add_flow_box(slide, title, body, left, top, width, height, accent):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.6)

    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.18), Inches(top + 0.2), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent
    dot.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(left + 0.46), Inches(top + 0.14), Inches(width - 0.6), Inches(0.3))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEAD
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT

    body_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.55), Inches(width - 0.4), Inches(height - 0.7))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.name = FONT_BODY
    run.font.size = Pt(12.5)
    run.font.color.rgb = MUTED


def add_arrow(slide, left, top, width):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(left), Inches(top), Inches(width), Inches(0.34))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SOFT_BLUE
    arrow.line.fill.background()


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_topbar(slide, slide_num=1)
    add_title_block(
        slide,
        "Prototype Update",
        "设备数据协同平台\n这一版原型进展",
        "这次重点想给大家看：界面已经更完整，也更像一套能顺着流程演示的系统了。",
        left=0.7,
        top=0.95,
        width=5.0,
    )
    add_chip(slide, "先提前说明：数据和设备名还是演示用", 0.72, 3.15, 3.35, SOFT_RED, BOSCH_RED_DARK)

    meta = slide.shapes.add_textbox(Inches(0.72), Inches(3.75), Inches(3.5), Inches(1.0))
    tf = meta.text_frame
    tf.clear()
    for idx, line in enumerate(["Bosch Power Tools · 杭州", "汇报人：YuYu", "日期：2026.04"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        run.font.name = FONT_BODY
        run.font.size = Pt(15)
        run.font.color.rgb = TEXT if idx == 0 else MUTED

    if AUTH_SCREEN.exists():
        add_image_panel(slide, AUTH_SCREEN, "当前登录页", 6.2, 0.75, 6.4, 5.55)
    elif AUTH_VISUAL.exists():
        add_image_panel(slide, AUTH_VISUAL, "Bosch 品牌视觉", 6.2, 0.75, 6.4, 5.55)
    add_footer(slide, "Bosch Power Tools | YuYu | Current UI Version")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_topbar(slide, slide_num=2)
    add_title_block(slide, "What changed", "这版我先做到了什么", "现在打开之后，已经更像一套能完整演示的系统，而不是几张分散页面。")
    add_card(slide, "登录页重做了", "现在的入口更完整，有品牌图、多语言切换和演示账号提示。", 0.72, 1.9, 2.85, 1.95, tag="01", accent=BOSCH_RED, fill=WHITE)
    add_card(slide, "管理端更顺了", "首页、趋势分析、设备看板、预测性维护这些主路径都能顺着点下去。", 3.85, 1.9, 2.85, 1.95, tag="02", accent=BOSCH_BLUE, fill=WHITE)
    add_card(slide, "不是只看页面了", "账号权限和审计日志也补进来了，系统味道更完整。", 6.98, 1.9, 2.85, 1.95, tag="03", accent=WARNING, fill=WHITE)
    add_card(slide, "后面好接真数据", "现在各个模块的位置基本定下来了，后面替换真实数据会更顺手。", 10.1, 1.9, 2.5, 1.95, tag="04", accent=SUCCESS, fill=WHITE)
    if ADMIN_SCREEN.exists():
        add_image_panel(slide, ADMIN_SCREEN, "当前首页", 0.78, 4.2, 7.15, 2.25)
    add_bullet_box(
        slide,
        "我想先让领导看到的，不是数字准不准",
        [
            "而是这个界面方向是不是对的。",
            "模块是不是够清楚。",
            "后面有没有继续往下做的价值。",
        ],
        8.18,
        4.18,
        4.4,
        2.28,
        fill=WHITE,
        accent=BOSCH_RED,
    )
    add_footer(slide, "Bosch Power Tools | YuYu | Current UI Version")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_topbar(slide, slide_num=3)
    add_title_block(slide, "Current UI", "现在界面基本长这样", "统一入口、顶部搜索和语言切换、左侧模块导航，这几个骨架已经稳定下来了。")
    if AUTH_SCREEN.exists():
        add_image_panel(slide, AUTH_SCREEN, "登录页：品牌图 + 多语言 + 演示账号", 0.72, 1.9, 5.35, 4.6)
    if ADMIN_SCREEN.exists():
        add_image_panel(slide, ADMIN_SCREEN, "管理首页：顶部全局区 + 左侧导航 + 产线总览", 6.28, 1.9, 6.25, 4.6)
    add_flow_box(slide, "统一登录", "管理员、编辑和只读账号都从同一个入口进。", 0.9, 6.7, 3.0, 0.9, BOSCH_RED)
    add_arrow(slide, 4.05, 7.0, 0.45)
    add_flow_box(slide, "角色识别", "登录后直接进入对应工作台，不用再切模式。", 4.65, 6.7, 2.9, 0.9, BOSCH_BLUE)
    add_arrow(slide, 7.75, 7.0, 0.45)
    add_flow_box(slide, "模块切换", "左边点模块，右边就能连续演示首页、趋势、设备、预测这些页面。", 8.35, 6.7, 4.2, 0.9, SUCCESS)
    add_footer(slide, "Bosch Power Tools | YuYu | Current UI Version")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_topbar(slide, slide_num=4)
    add_title_block(slide, "Key pages", "我觉得现在最值得看的几个页面", "这三页最能说明现在这个原型已经不只是“好看”，而是开始有逻辑了。")
    if ANALYTICS_SCREEN.exists():
        add_image_panel(slide, ANALYTICS_SCREEN, "趋势分析：OEE、T-loss、MTBF、MTTR 和停机原因", 0.72, 1.9, 4.0, 4.3)
    if EQUIPMENT_SCREEN.exists():
        add_image_panel(slide, EQUIPMENT_SCREEN, "设备看板：按阶段看设备状态，也能直接看到高风险项", 4.86, 1.9, 4.0, 4.3)
    if PREDICTIVE_SCREEN.exists():
        add_image_panel(slide, PREDICTIVE_SCREEN, "预测性维护：建议页已经能顺着风险、窗口和动作往下讲", 9.0, 1.9, 3.6, 4.3)
    add_card(
        slide,
        "为什么我会重点讲这三页",
        "因为它们基本对应了三个问题：\n1. 现在线上表现怎么样\n2. 哪台设备最该看\n3. 如果要提前处理，建议怎么做",
        0.78,
        6.42,
        11.8,
        0.95,
        tag="Main story",
        accent=BOSCH_RED,
        fill=WHITE,
    )
    add_footer(slide, "Bosch Power Tools | YuYu | Current UI Version")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_topbar(slide, slide_num=5)
    add_title_block(slide, "Collaboration", "这次不是只做展示页，也把协同这块补上了", "账号权限和审计日志补进去之后，这套原型看起来就更完整了。")
    if ACCOUNTS_SCREEN.exists():
        add_image_panel(slide, ACCOUNTS_SCREEN, "账号与权限：能创建普通账号，也能切编辑 / 只读", 0.72, 1.9, 6.0, 4.45)
    if AUDIT_SCREEN.exists():
        add_image_panel(slide, AUDIT_SCREEN, "审计日志：谁创建、谁提交、谁改了什么，都能往回看", 6.88, 1.9, 5.7, 4.45)
    add_bullet_box(
        slide,
        "另外，普通账号侧也已经有基本路径了",
        [
            "设备总览",
            "数据提交",
            "我的记录",
            "也就是说，后面不只是管理层看页面，业务侧也有机会真的用起来。",
        ],
        0.78,
        6.52,
        11.8,
        0.92,
        fill=SOFT_GRAY,
        accent=BOSCH_BLUE,
    )
    add_footer(slide, "Bosch Power Tools | YuYu | Current UI Version")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_topbar(slide, slide_num=6)
    add_title_block(slide, "Why it matters", "我觉得现在最有价值的几点", "如果这次是看“进展”，我觉得下面这四点最值得讲。")
    add_card(slide, "1. 界面骨架已经稳定了", "登录、顶部区、左侧导航、内容区这些核心骨架已经基本成型。", 0.78, 2.0, 3.0, 1.85, tag="01", accent=BOSCH_RED, fill=WHITE)
    add_card(slide, "2. 管理视角比较清楚", "首页、趋势、设备、预测这条主线已经能顺着讲清楚。", 3.95, 2.0, 3.0, 1.85, tag="02", accent=BOSCH_BLUE, fill=WHITE)
    add_card(slide, "3. 协同意识补进来了", "账号权限、审计日志、普通账号入口这些东西，让它不再只是静态展示页。", 7.12, 2.0, 3.0, 1.85, tag="03", accent=WARNING, fill=WHITE)
    add_card(slide, "4. 后面接真数据有落点", "现在每个模块放什么、怎么切换，已经比较清楚了，下一步更适合接真实内容。", 10.28, 2.0, 2.35, 1.85, tag="04", accent=SUCCESS, fill=WHITE)

    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.3), Inches(11.78), Inches(1.7))
    quote.fill.solid()
    quote.fill.fore_color.rgb = WHITE
    quote.line.color.rgb = LINE
    tf = quote.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "现在这版最想证明的一件事是：这个方向已经能“看起来像一套系统”，而不是停留在想法上。"
    run.font.name = FONT_HEAD
    run.font.size = Pt(21)
    run.font.bold = True
    run.font.color.rgb = TEXT
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = "等真实数据进来以后，这个原型会比从头再搭一版推进得更快。"
    run.font.name = FONT_BODY
    run.font.size = Pt(13)
    run.font.color.rgb = MUTED
    add_footer(slide, "Bosch Power Tools | YuYu | Current UI Version")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_topbar(slide, slide_num=7)
    add_title_block(slide, "Heads-up", "这里我想先提前说明一下", "这次更适合看方向和界面，不太适合把里面的数字当成正式结论。")
    add_bullet_box(
        slide,
        "先别太当真的部分",
        [
            "页面里的数据还是演示数据。",
            "设备名、风险项和部分记录也是为了演示效果临时放进去的。",
            "阈值和预测逻辑还没有跟业务一起细调。",
        ],
        0.75,
        1.95,
        5.7,
        3.0,
        fill=SOFT_RED,
        accent=BOSCH_RED,
    )
    add_bullet_box(
        slide,
        "这次更建议大家重点看",
        [
            "界面结构是不是顺。",
            "模块分法是不是合理。",
            "领导真正想看的内容，有没有在页面上露出来。",
        ],
        6.6,
        1.95,
        5.95,
        3.0,
        fill=SOFT_BLUE,
        accent=BOSCH_BLUE,
    )
    add_card(
        slide,
        "一句话概括",
        "这次先看“方向对不对”，先别纠结“数字真不真”。如果方向认可，下一步我就把重点切到真实数据和真实口径上。",
        0.82,
        5.35,
        11.7,
        1.25,
        tag="Positioning",
        accent=BOSCH_RED,
        fill=WHITE,
    )
    add_footer(slide, "Bosch Power Tools | YuYu | Current UI Version")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_topbar(slide, slide_num=8)
    add_title_block(slide, "Next step", "如果继续往下做，我建议这样推进", "我会把重点从“把页面做出来”，切到“把真实内容放进来”。")
    add_bullet_box(
        slide,
        "我建议的推进顺序",
        [
            "先确认现在这套界面方向是不是认可。",
            "再拿第一批真实设备清单和样例数据。",
            "然后把字段口径、阈值和规则一起对齐。",
            "最后做一版更像真的第二版原型。",
        ],
        0.75,
        1.85,
        5.55,
        3.25,
        fill=WHITE,
        accent=BOSCH_RED,
    )
    add_bullet_box(
        slide,
        "如果这次汇报后能拿到这几样，推进会快很多",
        [
            "一个业务接口人",
            "一批样例数据",
            "真实设备清单",
            "关键指标和阈值的初版口径",
        ],
        6.55,
        1.85,
        6.0,
        3.25,
        fill=WHITE,
        accent=BOSCH_BLUE,
    )
    add_card(
        slide,
        "我希望这次汇报结束后，至少先明确一件事",
        "这个方向值不值得继续往下做。如果答案是“值得”，那我下一步就优先把真实数据和真实口径接进来。",
        0.78,
        5.45,
        11.75,
        1.1,
        tag="My ask",
        accent=BOSCH_RED,
        fill=WHITE,
    )
    add_footer(slide, "Bosch Power Tools | YuYu | Current UI Version")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    return OUTPUT_FILE


if __name__ == "__main__":
    output = build_presentation()
    print(output)
