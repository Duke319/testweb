from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "artifacts" / "outputs"
OUTPUT_FILE = OUTPUT_DIR / "bosch-frontend-interface-jd-review.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BOSCH_RED = RGBColor(0xE2, 0x00, 0x15)
BOSCH_RED_DARK = RGBColor(0xA9, 0x00, 0x11)
BOSCH_BLUE = RGBColor(0x00, 0x56, 0x91)
CHARCOAL = RGBColor(0x1A, 0x1C, 0x1C)
MUTED = RGBColor(0x5E, 0x66, 0x73)
LINE = RGBColor(0xD6, 0xD9, 0xDE)
BG = RGBColor(0xF4, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_RED = RGBColor(0xFF, 0xEB, 0xE9)
SOFT_BLUE = RGBColor(0xE8, 0xF2, 0xFB)
SOFT_GRAY = RGBColor(0xF8, 0xF8, 0xF8)
SUCCESS = RGBColor(0x15, 0x80, 0x3D)
WARNING = RGBColor(0xD9, 0x77, 0x06)

FONT = "Microsoft YaHei"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_lines(slide, lines, x, y, w, h, size=13, color=CHARCOAL):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_top(slide, page=None, dark=False):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BOSCH_RED
    bar.line.fill.background()
    add_text(
        slide,
        "Bosch Power Tools · Frontend Direction Review",
        0.48,
        0.28,
        5.2,
        0.24,
        size=9.2,
        color=WHITE if dark else BOSCH_RED,
        bold=True,
    )
    if page:
        add_text(slide, f"{page:02d}", 12.34, 7.16, 0.44, 0.18, size=9, color=WHITE if dark else MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, dark=False):
    add_text(slide, "Bosch Power Tools | YuYu | 前端界面方向评审", 0.48, 7.16, 7.1, 0.18, size=8.3, color=WHITE if dark else MUTED)


def add_title(slide, kicker, title, subtitle="", page=None):
    add_top(slide, page=page)
    add_text(slide, kicker.upper(), 0.72, 0.72, 5.2, 0.24, size=9.3, color=BOSCH_RED, bold=True)
    add_text(slide, title, 0.72, 1.05, 7.5, 0.75, size=28, color=CHARCOAL, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.88, 8.2, 0.34, size=12.2, color=MUTED)
    add_footer(slide)


def add_pill(slide, text, x, y, w, fill, color=CHARCOAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(9.5)
    run.font.bold = True
    run.font.color.rgb = color


def add_card(slide, title, body, x, y, w, h, accent=BOSCH_RED, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    add_text(slide, title, x + 0.22, y + 0.2, w - 0.44, 0.28, size=15.2, color=CHARCOAL, bold=True)
    add_lines(slide, body.split("\n"), x + 0.22, y + 0.68, w - 0.44, h - 0.82, size=11.2, color=MUTED)


def add_step(slide, num, title, body, x, y, w, accent):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.48), Inches(0.48))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text(slide, num, x + 0.07, y + 0.11, 0.34, 0.18, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.64, y - 0.02, w - 0.64, 0.26, size=14.5, color=CHARCOAL, bold=True)
    add_text(slide, body, x + 0.64, y + 0.35, w - 0.64, 0.4, size=10.8, color=MUTED)


def add_flow(slide, items, x, y, w, accent=BOSCH_RED):
    gap = 0.18
    item_w = (w - gap * (len(items) - 1)) / len(items)
    for idx, item in enumerate(items):
        bx = x + idx * (item_w + gap)
        fill = SOFT_RED if idx == 0 else SOFT_BLUE if idx == 1 else SOFT_GRAY
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(y), Inches(item_w), Inches(0.78))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = LINE
        add_text(slide, item, bx + 0.1, y + 0.2, item_w - 0.2, 0.24, size=11.2, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(items) - 1:
            add_text(slide, ">", bx + item_w + 0.04, y + 0.24, 0.1, 0.2, size=13, color=accent, bold=True, align=PP_ALIGN.CENTER)


def add_wireframe(slide, x, y, w, h, variant="desktop"):
    outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    outer.fill.solid()
    outer.fill.fore_color.rgb = WHITE
    outer.line.color.rgb = LINE

    if variant == "login":
        left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.22), Inches(y + 0.24), Inches(w * 0.34), Inches(h - 0.48))
        left.fill.solid()
        left.fill.fore_color.rgb = SOFT_GRAY
        left.line.color.rgb = LINE
        right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + w * 0.42), Inches(y + 0.24), Inches(w * 0.48), Inches(h - 0.48))
        right.fill.solid()
        right.fill.fore_color.rgb = SOFT_BLUE
        right.line.fill.background()
        for idx in range(3):
            row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.42), Inches(y + 0.72 + idx * 0.48), Inches(w * 0.22), Inches(0.12))
            row.fill.solid()
            row.fill.fore_color.rgb = LINE if idx else BOSCH_RED
            row.line.fill.background()
        add_text(slide, "入口页", x + 0.52, y + h - 0.42, 0.9, 0.16, size=8.5, color=MUTED)
        return

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.18), Inches(y + 0.18), Inches(w - 0.36), Inches(0.42))
    top.fill.solid()
    top.fill.fore_color.rgb = SOFT_GRAY
    top.line.color.rgb = LINE

    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.18), Inches(y + 0.76), Inches(w * 0.22), Inches(h - 0.94))
    side.fill.solid()
    side.fill.fore_color.rgb = WHITE
    side.line.color.rgb = LINE

    content_x = x + 0.32 + w * 0.22
    content_w = w - 0.5 - w * 0.22
    hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(content_x), Inches(y + 0.76), Inches(content_w), Inches(0.58))
    hero.fill.solid()
    hero.fill.fore_color.rgb = SOFT_RED
    hero.line.color.rgb = LINE

    card_w = (content_w - 0.24) / 3
    for idx in range(3):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(content_x + idx * (card_w + 0.12)), Inches(y + 1.54), Inches(card_w), Inches(0.62))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LINE
    table = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(content_x), Inches(y + 2.36), Inches(content_w), Inches(h - 2.62))
    table.fill.solid()
    table.fill.fore_color.rgb = SOFT_GRAY
    table.line.color.rgb = LINE

    for idx in range(6):
        nav = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.36), Inches(y + 1.03 + idx * 0.32), Inches(w * 0.14), Inches(0.1))
        nav.fill.solid()
        nav.fill.fore_color.rgb = BOSCH_RED if idx == 0 else LINE
        nav.line.fill.background()


def add_module_map(slide, x, y):
    modules = [
        ("总览", "先看到整体状态"),
        ("分析", "理解趋势变化"),
        ("设备", "定位重点对象"),
        ("规则", "解释风险依据"),
        ("扩展", "保留未来想法"),
    ]
    for idx, (title, body) in enumerate(modules):
        col = idx % 3
        row = idx // 3
        bx = x + col * 2.55
        by = y + row * 1.32
        fill = SOFT_RED if idx == 0 else SOFT_BLUE if idx in (1, 2) else SOFT_GRAY
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(by), Inches(2.18), Inches(0.95))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = LINE
        add_text(slide, title, bx + 0.16, by + 0.16, 1.85, 0.22, size=13.5, color=CHARCOAL, bold=True)
        add_text(slide, body, bx + 0.16, by + 0.5, 1.85, 0.22, size=9.8, color=MUTED)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.author = "YuYu"
    prs.core_properties.title = "设备管理前端界面方向评审"

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CHARCOAL)
    add_top(slide, page=1, dark=True)
    add_text(slide, "FRONTEND DIRECTION", 0.72, 0.88, 4.2, 0.24, size=9.5, color=WHITE, bold=True)
    add_text(slide, "设备管理\n前端界面方向\n阶段汇报", 0.72, 1.35, 5.2, 1.75, size=32, color=WHITE, bold=True)
    add_text(slide, "基于实习岗位 JD 先做出的页面方向假设", 0.74, 3.36, 5.1, 0.32, size=13, color=RGBColor(0xEA, 0xEA, 0xEA))
    add_pill(slide, "本次只看前端方向", 0.74, 4.05, 1.95, SOFT_RED, BOSCH_RED_DARK)
    add_lines(slide, ["场景：设备管理数字化探索", "汇报人：YuYu", "日期：2026.04"], 0.76, 4.62, 3.7, 0.82, size=12, color=RGBColor(0xEA, 0xEA, 0xEA))
    add_wireframe(slide, 6.35, 1.0, 5.5, 4.85, variant="desktop")
    add_text(slide, "不是最终系统演示，而是页面方向评审", 6.65, 6.08, 4.9, 0.28, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, dark=True)

    # 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "Stage 1 / Starting Point", "为什么先做这个前端界面", "先把基于 JD 的个人理解可视化，再请 Boss 判断方向是否一致。", page=2)
    add_step(slide, "01", "输入来自实习岗位 JD", "先根据岗位描述理解可能涉及的设备管理、数据分析和数字化方向。", 0.92, 2.45, 3.8, BOSCH_RED)
    add_step(slide, "02", "业务信息还不完整", "目前还没有完整了解真实产线、设备清单、字段口径和现场数据。", 4.95, 2.45, 3.85, BOSCH_BLUE)
    add_step(slide, "03", "先做页面方向假设", "用前端页面表达我对管理视角和信息层级的理解。", 9.0, 2.45, 3.45, WARNING)
    add_flow(slide, ["实习岗位 JD", "我的理解", "前端方向", "Boss 反馈", "页面调整"], 0.96, 4.5, 11.25)
    add_card(slide, "这一页的核心", "我不是直接给出最终方案，而是先用页面结构表达理解，让 Boss 更容易判断方向是否走对。", 1.02, 5.62, 11.15, 1.05, accent=BOSCH_RED, fill=WHITE)

    # 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "Stage 1 / Scope", "这次控制展示范围", "只看页面方向、信息层级和视觉结构，不展开系统实现。", page=3)
    add_card(slide, "本次主要展示", "页面入口的整体感\n导航和内容区域的关系\n管理视角的信息层级\n普通账号页面是否有必要", 0.92, 2.32, 3.6, 3.15, accent=BOSCH_RED, fill=SOFT_RED)
    add_card(slide, "本次先不展开", "真实数据准确性\n产线和设备明细\n流程是否完整落地\n系统实现和交互细节", 4.92, 2.32, 3.6, 3.15, accent=BOSCH_BLUE, fill=SOFT_BLUE)
    add_card(slide, "希望得到反馈", "页面方向是否匹配期待\n首页优先级是否合理\n哪些页面需要保留或删减\n下一版先调整哪里", 8.92, 2.32, 3.5, 3.15, accent=SUCCESS, fill=SOFT_GRAY)
    add_pill(slide, "关键词：方向先行，细节后补", 4.64, 6.05, 3.25, SOFT_RED, BOSCH_RED_DARK)

    # 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "Stage 2 / Page Skeleton", "我先搭了一个页面骨架", "用抽象结构表达页面关系，不在这一页展示具体界面细节。", page=4)
    add_wireframe(slide, 0.9, 2.12, 4.1, 3.82, variant="login")
    add_wireframe(slide, 5.48, 2.12, 6.65, 3.82, variant="desktop")
    add_flow(slide, ["入口页", "顶部全局区", "左侧导航", "内容展示区"], 1.32, 6.28, 10.48)
    add_text(slide, "请 Boss 看：这样的页面骨架是否符合设备管理场景的第一印象。", 1.2, 6.95, 10.9, 0.26, size=12.8, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

    # 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "Stage 2 / Management View", "管理员视角先按五类页面组织", "不展示过细模块清单，只保留页面方向和信息分组。", page=5)
    add_module_map(slide, 0.92, 2.32)
    add_card(slide, "我这样分组的原因", "管理视角通常先需要看到整体状态，再继续理解趋势、定位设备、解释风险。预测或扩展页面先作为方向保留。", 0.92, 5.12, 6.95, 0.95, accent=BOSCH_RED, fill=SOFT_RED)
    add_card(slide, "需要 Boss 判断", "这五类页面是否符合岗位期待？是否应该减少页面数量，先集中展示最核心的管理视角？", 8.25, 2.32, 3.88, 3.75, accent=BOSCH_BLUE, fill=SOFT_BLUE)

    # 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "Stage 2 / Information Flow", "设备管理页面的展示主线", "我的假设：页面应该帮助管理层从整体状态走向重点对象。", page=6)
    add_flow(slide, ["整体状态", "趋势变化", "重点设备", "风险解释"], 1.0, 2.42, 11.15)
    add_card(slide, "整体状态", "首页只承担方向入口：让人先知道设备状态大概如何。", 0.92, 3.82, 2.85, 1.55, accent=BOSCH_RED, fill=WHITE)
    add_card(slide, "趋势变化", "用图表区域表达效率或风险随时间变化的可能呈现方式。", 3.92, 3.82, 2.85, 1.55, accent=BOSCH_BLUE, fill=WHITE)
    add_card(slide, "重点设备", "用卡片或列表把注意力落到具体设备对象上。", 6.92, 3.82, 2.85, 1.55, accent=WARNING, fill=WHITE)
    add_card(slide, "风险解释", "用简单规则区域解释为什么这些对象需要关注。", 9.92, 3.82, 2.42, 1.55, accent=SUCCESS, fill=WHITE)
    add_text(slide, "请 Boss 判断：这条观看顺序是否符合管理层真正的判断习惯。", 1.2, 6.24, 10.9, 0.28, size=13, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

    # 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "Stage 2 / Role Assumption", "普通账号侧只作为一个页面假设", "我不确定现场侧是否需要单独页面，所以先把它作为待判断方向。", page=7)
    add_card(slide, "如果需要普通账号侧", "可以展示设备状态、填写变化信息、查看自己的记录。", 0.92, 2.34, 3.58, 2.65, accent=BOSCH_RED, fill=SOFT_RED)
    add_card(slide, "如果不需要普通账号侧", "可以把汇报重点收回管理员视角，让页面更聚焦。", 4.86, 2.34, 3.58, 2.65, accent=BOSCH_BLUE, fill=SOFT_BLUE)
    add_card(slide, "这页只问一个问题", "普通账号侧是否符合真实工作方式？如果不符合，下一版直接弱化或删除。", 8.78, 2.34, 3.42, 2.65, accent=SUCCESS, fill=SOFT_GRAY)
    add_flow(slide, ["查看状态", "填写信息", "查看记录"], 1.38, 5.72, 10.25, accent=BOSCH_BLUE)

    # 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CHARCOAL)
    add_top(slide, page=8, dark=True)
    add_text(slide, "STAGE 3 / FEEDBACK", 0.72, 0.88, 4.2, 0.24, size=9.5, color=WHITE, bold=True)
    add_text(slide, "这次先确认\n前端页面方向", 0.72, 1.32, 4.95, 1.25, size=31, color=WHITE, bold=True)
    add_text(slide, "如果方向一致，下一版再根据反馈调整页面结构和展示重点。", 0.74, 2.85, 4.9, 0.42, size=13, color=RGBColor(0xEA, 0xEA, 0xEA))
    add_card(slide, "希望 Boss 判断", "1. 页面方向是否符合岗位期待\n2. 管理视角是否足够聚焦\n3. 普通账号侧是否有必要\n4. 页面主线是否符合观看习惯", 6.15, 0.98, 5.85, 2.32, accent=BOSCH_RED, fill=WHITE)
    add_card(slide, "下一步建议", "1. 先调整页面骨架和模块取舍\n2. 强化最重要的管理视角\n3. 弱化不必要的页面细节\n4. 第二版再补更贴近业务的内容", 6.15, 3.65, 5.85, 2.1, accent=BOSCH_BLUE, fill=WHITE)
    add_flow(slide, ["方向确认", "页面调整", "模块取舍", "第二版界面"], 0.74, 4.25, 5.15)
    add_footer(slide, dark=True)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    return OUTPUT_FILE


if __name__ == "__main__":
    print(build())
